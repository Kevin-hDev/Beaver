import hashlib
import tempfile
from pathlib import Path
from urllib.request import Request, urlopen

import xlsxwriter
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas


FIXTURES = Path(__file__).resolve().parent
FONT_URL = (
    "https://raw.githubusercontent.com/notofonts/noto-cjk/"
    "f8d157532fbfaeda587e826d4cd5b21a49186f7c/"
    "Sans/Variable/TTF/Subset/NotoSansJP-VF.ttf"
)
FONT_SHA256 = "f4b373b226668ee33a6e54b02823dcd2d1209f17159f777421ae8c2275160369"
MAX_FONT_BYTES = 10 * 1024 * 1024


def create_spreadsheet() -> None:
    target = FIXTURES / "external-large-preview.xlsx"
    workbook = xlsxwriter.Workbook(target, {"constant_memory": False})
    sheet = workbook.add_worksheet("External")
    shared_value = "界" * 32_767
    for row in range(200):
        for column in range(100):
            sheet.write_string(row, column, shared_value)
    workbook.close()


def create_pdf() -> None:
    target = FIXTURES / "external-large-text.pdf"
    with tempfile.TemporaryDirectory(prefix="beaver-fixture-font-") as directory:
        font_path = Path(directory) / "NotoSansJP-VF.ttf"
        download_font(font_path)
        pdfmetrics.registerFont(TTFont("ExternalNotoSansJP", font_path))
        document = canvas.Canvas(str(target), pageCompression=1)
        document.setFont("ExternalNotoSansJP", 4)
        line = "日本語の外部PDFテキスト" * 12
        for _page in range(50):
            y = 820
            for _line in range(100):
                document.drawString(24, y, line)
                y -= 8
            document.showPage()
            document.setFont("ExternalNotoSansJP", 4)
        document.save()


def download_font(target: Path) -> None:
    request = Request(FONT_URL, headers={"User-Agent": "Beaver fixture generator"})
    digest = hashlib.sha256()
    total = 0
    with urlopen(request, timeout=30) as response, target.open("wb") as output:
        length = int(response.headers.get("Content-Length", "0"))
        if length <= 0 or length > MAX_FONT_BYTES:
            raise ValueError("invalid font size")
        while chunk := response.read(64 * 1024):
            total += len(chunk)
            if total > MAX_FONT_BYTES:
                raise ValueError("font exceeds size limit")
            digest.update(chunk)
            output.write(chunk)
    if total != length or digest.hexdigest() != FONT_SHA256:
        raise ValueError("font checksum mismatch")


def create_presentation() -> None:
    target = FIXTURES / "external-fragmented-runs.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    paragraph = box.text_frame.paragraphs[0]
    fragments = [
        ("Prepared for {{", False),
        ("cus", True),
        ("tom", False),
        ("er}}", True),
    ]
    for text, bold in fragments:
        run = paragraph.add_run()
        run.text = text
        run.font.bold = bold
        run.font.size = Pt(24)
    presentation.save(target)


def create_unicode_sample() -> None:
    target = FIXTURES / "external-unicode-samples.txt"
    target.write_text(
        "日本語 中文 한국어 Русский العربية",
        encoding="utf-8",
    )


if __name__ == "__main__":
    create_spreadsheet()
    create_pdf()
    create_presentation()
    create_unicode_sample()
